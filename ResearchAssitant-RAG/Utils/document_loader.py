from __future__ import annotations

import json
import mimetypes
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple, Union


@dataclass(frozen=True)
class LoadedDocument:
    """
    Minimal, RAG-friendly document container.

    - `text`: extracted plain text content
    - `metadata`: source info (path, type, page/sheet, etc.)
    """

    text: str
    metadata: Dict[str, Any]


PathLike = Union[str, Path]


def load_documents(
    inputs: Union[PathLike, Sequence[PathLike]],
    *,
    recursive: bool = True,
    ignore_hidden: bool = True,
    allowed_extensions: Optional[Sequence[str]] = None,
) -> List[LoadedDocument]:
    """
    Load one file, many files, or a directory into a list of `LoadedDocument`.

    Supported (best-effort) formats by extension:
    - .txt, .md
    - .pdf (requires `pypdf`)
    - .docx (requires `python-docx`)
    - .pptx (requires `python-pptx`)
    - .html/.htm (requires `beautifulsoup4`)
    - .csv (built-in csv)
    - .xlsx/.xlsm (requires `openpyxl`)
    - .json (built-in json)

    Notes:
    - Some formats require optional dependencies; if missing, a clear ImportError is raised.
    - Returns multiple documents for multi-part sources (e.g., PDFs per page, Excel per sheet).
    """

    paths = _normalize_inputs(inputs)
    file_paths: List[Path] = []
    for p in paths:
        if p.is_dir():
            file_paths.extend(_iter_files(p, recursive=recursive, ignore_hidden=ignore_hidden))
        else:
            file_paths.append(p)

    if allowed_extensions is not None:
        allowed = {e.lower().lstrip(".") for e in allowed_extensions}
        file_paths = [p for p in file_paths if p.suffix.lower().lstrip(".") in allowed]

    docs: List[LoadedDocument] = []
    for path in file_paths:
        if ignore_hidden and _is_hidden(path):
            continue
        if not path.exists() or not path.is_file():
            continue
        docs.extend(load_document(path))

    return docs


def load_document(path: PathLike) -> List[LoadedDocument]:
    """Load a single file into one-or-more `LoadedDocument` entries."""
    p = Path(path).expanduser().resolve()
    ext = p.suffix.lower()

    if ext in {".txt", ".md"}:
        return [_load_text_file(p)]
    if ext == ".pdf":
        return _load_pdf(p)
    if ext == ".docx":
        return [_load_docx(p)]
    if ext == ".pptx":
        return _load_pptx(p)
    if ext in {".html", ".htm"}:
        return [_load_html(p)]
    if ext == ".csv":
        return [_load_csv(p)]
    if ext in {".xlsx", ".xlsm"}:
        return _load_xlsx(p)
    if ext == ".json":
        return [_load_json(p)]

    # Fallback: try to read as text
    return [_load_text_file(p)]


def _normalize_inputs(inputs: Union[PathLike, Sequence[PathLike]]) -> List[Path]:
    if isinstance(inputs, (str, Path)):
        inputs_list: Sequence[PathLike] = [inputs]
    else:
        inputs_list = inputs
    return [Path(p).expanduser().resolve() for p in inputs_list]


def _iter_files(root: Path, *, recursive: bool, ignore_hidden: bool) -> Iterable[Path]:
    it = root.rglob("*") if recursive else root.glob("*")
    for p in it:
        if p.is_dir():
            continue
        if ignore_hidden and _is_hidden(p):
            continue
        yield p


def _is_hidden(p: Path) -> bool:
    # Cross-platform heuristic: dotfiles/dirs are "hidden" everywhere.
    # Windows hidden attribute isn't checked to avoid platform-specific APIs.
    return any(part.startswith(".") for part in p.parts)


def _base_metadata(path: Path) -> Dict[str, Any]:
    mime, _ = mimetypes.guess_type(str(path))
    return {
        "source": str(path),
        "filename": path.name,
        "extension": path.suffix.lower().lstrip("."),
        "mime_type": mime,
        "loaded_at": datetime.now(timezone.utc).isoformat(),
    }


def _read_text_best_effort(path: Path) -> str:
    data = path.read_bytes()
    # UTF-8 (with BOM) then fallback.
    for enc in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _load_text_file(path: Path) -> LoadedDocument:
    text = _read_text_best_effort(path)
    return LoadedDocument(text=text, metadata=_base_metadata(path))


def _load_pdf(path: Path) -> List[LoadedDocument]:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception as e:  # pragma: no cover
        raise ImportError(
            "PDF loading requires `pypdf`. Install with: pip install pypdf"
        ) from e

    reader = PdfReader(str(path))
    docs: List[LoadedDocument] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        md = _base_metadata(path)
        md.update({"page": i, "page_count": len(reader.pages), "format": "pdf"})
        docs.append(LoadedDocument(text=text, metadata=md))
    return docs


def _load_docx(path: Path) -> LoadedDocument:
    try:
        import docx  # type: ignore
    except Exception as e:  # pragma: no cover
        raise ImportError(
            "DOCX loading requires `python-docx`. Install with: pip install python-docx"
        ) from e

    d = docx.Document(str(path))
    parts: List[str] = []
    for para in d.paragraphs:
        t = (para.text or "").strip()
        if t:
            parts.append(t)
    text = "\n".join(parts)
    md = _base_metadata(path)
    md.update({"format": "docx"})
    return LoadedDocument(text=text, metadata=md)


def _load_pptx(path: Path) -> List[LoadedDocument]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:  # pragma: no cover
        raise ImportError(
            "PPTX loading requires `python-pptx`. Install with: pip install python-pptx"
        ) from e

    prs = Presentation(str(path))
    docs: List[LoadedDocument] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str):
                t = text.strip()
                if t:
                    parts.append(t)
        md = _base_metadata(path)
        md.update({"format": "pptx", "slide": slide_idx, "slide_count": len(prs.slides)})
        docs.append(LoadedDocument(text="\n".join(parts), metadata=md))
    return docs


def _load_html(path: Path) -> LoadedDocument:
    try:
        from bs4 import BeautifulSoup  # type: ignore
    except Exception as e:  # pragma: no cover
        raise ImportError(
            "HTML loading requires `beautifulsoup4`. Install with: pip install beautifulsoup4"
        ) from e

    html = _read_text_best_effort(path)
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    md = _base_metadata(path)
    md.update({"format": "html"})
    return LoadedDocument(text=text, metadata=md)


def _load_csv(path: Path) -> LoadedDocument:
    import csv

    text = _read_text_best_effort(path)
    # Normalize as tab-separated "rows" to make it LLM-friendly while preserving structure.
    rows: List[str] = []
    reader = csv.reader(text.splitlines())
    for row in reader:
        rows.append("\t".join(str(cell) for cell in row))
    md = _base_metadata(path)
    md.update({"format": "csv"})
    return LoadedDocument(text="\n".join(rows), metadata=md)


def _load_xlsx(path: Path) -> List[LoadedDocument]:
    try:
        import openpyxl  # type: ignore
    except Exception as e:  # pragma: no cover
        raise ImportError(
            "Excel loading requires `openpyxl`. Install with: pip install openpyxl"
        ) from e

    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    docs: List[LoadedDocument] = []
    for ws in wb.worksheets:
        lines: List[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            # Trim trailing empties for cleaner output
            while cells and cells[-1] == "":
                cells.pop()
            if cells:
                lines.append("\t".join(cells))
        md = _base_metadata(path)
        md.update({"format": "xlsx", "sheet": ws.title})
        docs.append(LoadedDocument(text="\n".join(lines), metadata=md))
    return docs


def _load_json(path: Path) -> LoadedDocument:
    raw = _read_text_best_effort(path)
    try:
        obj = json.loads(raw)
        text = json.dumps(obj, ensure_ascii=False, indent=2)
    except json.JSONDecodeError:
        # Not valid JSON; keep raw for visibility
        text = raw
    md = _base_metadata(path)
    md.update({"format": "json"})
    return LoadedDocument(text=text, metadata=md)

